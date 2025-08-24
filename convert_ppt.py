from pptx import Presentation
import sys

def convert_pptx_to_md(pptx_file, md_file):
    prs = Presentation(pptx_file)
    content = [f"# {prs.core_properties.title or 'Presentation'}\n"]

    for i, slide in enumerate(prs.slides):
        content.append(f"\n## Slide {i+1}\n")
        for shape in slide.shapes:
            if shape.has_text_frame and shape.text.strip():
                content.append(f"{shape.text}\n")
        content.append("\n---\n")

    with open(md_file, 'w') as f:
        f.write(''.join(content))

    print(f"✅ Converted {pptx_file} to {md_file}")

if __name__ == "__main__":
    convert_pptx_to_md(sys.argv[1], sys.argv[2])
