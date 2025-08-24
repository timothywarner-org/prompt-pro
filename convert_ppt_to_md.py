#!/usr/bin/env python3
"""
PowerPoint to Markdown Converter
Converts PowerPoint presentations to Markdown format
"""

from pptx import Presentation
import sys
import os

def extract_text_from_shape(shape):
    """Extract text from a shape"""
    if hasattr(shape, "text"):
        return shape.text.strip()
    return ""

def extract_text_from_table(table):
    """Extract text from a table"""
    text = []
    for row in table.rows:
        row_text = []
        for cell in row.cells:
            if cell.text.strip():
                row_text.append(cell.text.strip())
        if row_text:
            text.append(" | ".join(row_text))
    return "\n".join(text)

def pptx_to_markdown(pptx_path, output_path=None):
    """Convert PowerPoint to Markdown"""

    if not os.path.exists(pptx_path):
        print(f"Error: File {pptx_path} not found")
        return False

    try:
        # Load the presentation
        prs = Presentation(pptx_path)

        # Start building markdown content
        markdown_content = []

        # Add title
        if prs.core_properties.title:
            markdown_content.append(f"# {prs.core_properties.title}\n")
        else:
            markdown_content.append("# Presentation\n")

        # Add metadata
        if prs.core_properties.author:
            markdown_content.append(f"**Author:** {prs.core_properties.author}\n")
        if prs.core_properties.created:
            markdown_content.append(f"**Created:** {prs.core_properties.created}\n")
        if prs.core_properties.modified:
            markdown_content.append(f"**Modified:** {prs.core_properties.modified}\n")

        markdown_content.append("\n---\n")

        # Process each slide
        for slide_num, slide in enumerate(prs.slides, 1):
            markdown_content.append(f"\n## Slide {slide_num}\n")

            # Extract slide title
            slide_title = ""
            for shape in slide.shapes:
                if shape.has_text_frame and shape.text.strip():
                    # Check if this looks like a title (usually the first text shape)
                    if not slide_title:
                        slide_title = shape.text.strip()
                        markdown_content.append(f"### {slide_title}\n")
                    else:
                        # Additional text content
                        text = shape.text.strip()
                        if text and text != slide_title:
                            markdown_content.append(f"{text}\n")

            # Process other shapes (tables, images, etc.)
            for shape in slide.shapes:
                # Handle tables
                if shape.has_table:
                    table_text = extract_text_from_table(shape.table)
                    if table_text:
                        markdown_content.append(f"\n**Table:**\n```\n{table_text}\n```\n")

                # Handle images (just note their presence)
                if shape.shape_type == 13:  # Picture
                    markdown_content.append("\n*[Image included in original presentation]*\n")

                # Handle charts
                if shape.has_chart:
                    markdown_content.append("\n*[Chart included in original presentation]*\n")

            markdown_content.append("\n---\n")

        # Write to file
        if output_path is None:
            output_path = pptx_path.replace('.pptx', '.md')

        with open(output_path, 'w', encoding='utf-8') as f:
            f.write(''.join(markdown_content))

        print(f"✅ Successfully converted {pptx_path} to {output_path}")
        return True

    except Exception as e:
        print(f"❌ Error converting presentation: {str(e)}")
        return False

if __name__ == "__main__":
    if len(sys.argv) < 2:
        print("Usage: python convert_ppt_to_md.py <pptx_file> [output_file]")
        sys.exit(1)

    pptx_file = sys.argv[1]
    output_file = sys.argv[2] if len(sys.argv) > 2 else None

    success = pptx_to_markdown(pptx_file, output_file)
    sys.exit(0 if success else 1)
